#!/usr/bin/env python3
"""Build a polished .pptx from slides.json.

Usage: build_pptx.py <talk_dir> [src.json] [dst.pptx]
"""
import json, io, os, sys, urllib.request, re, unicodedata
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

if len(sys.argv) < 2:
    print("usage: build_pptx.py <talk_dir> [src.json] [dst.pptx]"); sys.exit(1)

def slugify(s):
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode()
    s = re.sub(r"[^\w\s-]", "", s).strip().lower()
    return re.sub(r"[-\s]+", "-", s) or "slides"

TALK = os.path.abspath(sys.argv[1])
SRC = sys.argv[2] if len(sys.argv) > 2 else os.path.join(TALK, "slides.json")
_doc_peek = json.load(open(SRC))
_title = (_doc_peek.get("presentation") or {}).get("title") or "slides"
_default_dst = os.path.join(TALK, "output", f"{slugify(_title)}.pptx")
DST = sys.argv[3] if len(sys.argv) > 3 else _default_dst
os.makedirs(os.path.dirname(DST), exist_ok=True)

doc = json.load(open(SRC))
theme = doc.get("theme", {}).get("config", {})
colors = theme.get("colors", {})

def hx(h, d="#FFFFFF"):
    h = (h or d).lstrip("#")
    if len(h) != 6: h = d.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG       = hx(colors.get("background"), "#000000")
PRIMARY  = hx(colors.get("primary"),    "#FFFFFF")
SECONDARY= hx(colors.get("secondary"),  "#FF4013")
TEXT     = hx(colors.get("text"),       "#FFFFFF")
MUTED    = RGBColor(0x88,0x88,0x88)
PANEL    = RGBColor(0x14,0x14,0x14)
BORDER   = RGBColor(0x33,0x33,0x33)
FONT_H   = theme.get("fonts", {}).get("heading", "Inter") or "Inter"
FONT_B   = theme.get("fonts", {}).get("body", "Inter") or "Inter"
FONT_C   = (theme.get("fonts", {}).get("code") or "Menlo").strip() or "Menlo"
TALK_TITLE = doc.get("presentation", {}).get("title", "")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_bg(s, color=BG):
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background(); bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg

def rect(s, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    r = s.shapes.add_shape(shape, l, t, w, h)
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line
    return r

def text(s, l, t, w, h, content, *, size=20, bold=False, color=TEXT,
         font=FONT_B, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run(); r.text = str(line)
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return tb

def bullet_list(s, l, t, w, h, items, *, size=22, color=TEXT, font=FONT_B,
                bullet_color=None, gap=Pt(10), line_spacing=1.15):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    bc = bullet_color or SECONDARY
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = gap; p.line_spacing = line_spacing
        r1 = p.add_run(); r1.text = "•  "
        r1.font.name = font; r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bc
        r2 = p.add_run(); r2.text = str(item)
        r2.font.name = font; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb

def accent_bar(s, l, t, w=Inches(1.4), h=Inches(0.14), color=SECONDARY):
    return rect(s, l, t, w, h, fill=color)

def fetch(url):
    try:
        req = urllib.request.Request(url, headers={"User-Agent":"Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10) as r: return io.BytesIO(r.read())
    except Exception as e:
        print(f"  [warn] {e}", file=sys.stderr); return None

def fit_picture(s, png_path, area_l, area_t, area_w, area_h):
    """Insert PNG keeping aspect ratio, centered in the given area."""
    try:
        with Image.open(png_path) as im: iw, ih = im.size
    except Exception:
        s.shapes.add_picture(png_path, area_l, area_t, width=area_w, height=area_h); return
    ratio = iw / ih
    aw_emu, ah_emu = area_w, area_h
    if aw_emu / ah_emu > ratio:
        # area wider than image -> match height
        h = ah_emu; w = int(h * ratio)
    else:
        w = aw_emu; h = int(w / ratio)
    left = area_l + (aw_emu - w) // 2
    top  = area_t + (ah_emu - h) // 2
    s.shapes.add_picture(png_path, left, top, width=w, height=h)

# ---------- chrome ----------
def slide_number_footer(s, idx, total, show_title=True):
    if show_title and TALK_TITLE:
        text(s, Inches(0.6), Inches(7.05), Inches(8.0), Inches(0.35),
             TALK_TITLE, size=10, color=MUTED, font=FONT_B)
    text(s, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35),
         f"{idx:02d} / {total:02d}", size=10, color=MUTED, font=FONT_B,
         align=PP_ALIGN.RIGHT)

# ---------- templates ----------
def t_cover(s, d):
    add_bg(s)
    # decorative grid of dots top-right
    for col in range(7):
        for row in range(5):
            cx = Inches(10.6 + col*0.35); cy = Inches(0.6 + row*0.35)
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.08), Inches(0.08))
            dot.line.fill.background(); dot.fill.solid()
            dot.fill.fore_color.rgb = MUTED if (col+row) % 2 else SECONDARY
    accent_bar(s, Inches(0.8), Inches(2.55), w=Inches(1.8), h=Inches(0.16))
    text(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.7),
         d.get("title",""), size=78, bold=True, color=PRIMARY, font=FONT_H)
    if d.get("subtitle"):
        text(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.9),
             d["subtitle"], size=28, color=TEXT, font=FONT_B)
    if d.get("author"):
        text(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
             d["author"], size=18, color=SECONDARY, font=FONT_B)

def t_bio(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.9),
         d.get("title") or d.get("name", ""), size=44, bold=True, color=PRIMARY, font=FONT_H)
    if d.get("role"):
        text(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.5),
             d["role"], size=22, color=SECONDARY, font=FONT_B)
    if d.get("photo_url"):
        data = fetch(d["photo_url"])
        if data:
            try: s.shapes.add_picture(data, Inches(0.8), Inches(2.8),
                                      width=Inches(3.6), height=Inches(3.6))
            except Exception as e: print(f"  [warn] picture {e}", file=sys.stderr)
    bx = Inches(4.9) if d.get("photo_url") else Inches(0.8)
    bw = Inches(7.5) if d.get("photo_url") else Inches(11.7)
    bullet_list(s, bx, Inches(2.9), bw, Inches(3.6),
                d.get("bullets") or [], size=24)
    # Contact line at bottom
    contacts = []
    gh = d.get("github") or d.get("github_username")
    if gh: contacts.append(f"GitHub: @{gh.lstrip('@')}")
    if d.get("linkedin"): contacts.append(f"LinkedIn: {d['linkedin']}")
    if d.get("twitter"): contacts.append(f"Twitter: @{d['twitter'].lstrip('@')}")
    if d.get("website"): contacts.append(d["website"])
    if contacts:
        text(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
             "  ·  ".join(contacts), size=18, color=MUTED, font=FONT_B)

def t_agenda(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.9),
         d.get("title","Agenda"), size=44, bold=True, color=PRIMARY, font=FONT_H)
    items = d.get("blocks") or []
    # two columns
    col_w = Inches(5.7); col_h = Inches(4.6)
    col1_x = Inches(0.8); col2_x = Inches(6.9); top = Inches(2.2)
    half = (len(items) + 1) // 2
    for col, sub in enumerate([items[:half], items[half:]]):
        tb = s.shapes.add_textbox(col1_x if col==0 else col2_x, top, col_w, col_h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, item in enumerate(sub):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(14); p.line_spacing = 1.2
            r = p.add_run(); r.text = item
            r.font.name = FONT_B; r.font.size = Pt(24); r.font.color.rgb = TEXT

def t_section(s, d):
    add_bg(s)
    block = d.get("block") or d.get("number") or d.get("section_number")
    if block:
        text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.6),
             f"BLOCO {block}", size=18, color=SECONDARY, font=FONT_B, bold=True)
    accent_bar(s, Inches(0.8), Inches(3.3), w=Inches(2.4), h=Inches(0.16))
    text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(2.0),
         d.get("title",""), size=66, bold=True, color=PRIMARY, font=FONT_H,
         anchor=MSO_ANCHOR.TOP)

def t_question(s, d):
    add_bg(s)
    # giant centered question, accent over
    accent_bar(s, Inches(6.07), Inches(2.6), w=Inches(1.2), h=Inches(0.16))
    text(s, Inches(0.6), Inches(2.95), Inches(12.1), Inches(2.2),
         d.get("question",""), size=54, bold=True, color=PRIMARY, font=FONT_H,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.15)

def t_answer(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.4),
         d.get("question",""), size=34, bold=True, color=PRIMARY, font=FONT_H,
         line_spacing=1.15)
    bullet_list(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(4.0),
                d.get("bullets") or [], size=26, gap=Pt(14))

def t_content(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.0),
         d.get("title",""), size=38, bold=True, color=PRIMARY, font=FONT_H)
    y = Inches(1.9)
    if d.get("subtitle"):
        text(s, Inches(0.8), y, Inches(11.7), Inches(0.5),
             d["subtitle"], size=20, color=SECONDARY, font=FONT_B); y = Inches(2.45)
    glossary = d.get("glossary") or []
    bullets_h = Inches(4.3) if not glossary else Inches(2.9)
    bullet_list(s, Inches(0.8), y, Inches(11.7), bullets_h,
                d.get("bullets") or [], size=24)
    if glossary:
        gy = Inches(5.6)
        rect(s, Inches(0.8), gy, Inches(11.7), Inches(1.4),
             fill=PANEL, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, Inches(1.0), gy + Inches(0.12), Inches(11.3), Inches(0.3),
             "GLOSSÁRIO", size=11, bold=True, color=SECONDARY, font=FONT_B)
        # 2 columns of definitions
        for i, item in enumerate(glossary[:4]):
            col = i % 2; row = i // 2
            tx = Inches(1.0) + col*Inches(5.85)
            ty = gy + Inches(0.45) + row*Inches(0.45)
            term = item.get("term", "")
            defi = item.get("def", "")
            tb = s.shapes.add_textbox(tx, ty, Inches(5.7), Inches(0.45))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            r1 = p.add_run(); r1.text = f"{term} "
            r1.font.name = FONT_B; r1.font.size = Pt(13); r1.font.bold = True
            r1.font.color.rgb = SECONDARY
            r2 = p.add_run(); r2.text = defi
            r2.font.name = FONT_B; r2.font.size = Pt(13); r2.font.color.rgb = TEXT
    if d.get("quote"):
        text(s, Inches(0.8), Inches(6.4) if not glossary else Inches(7.05),
             Inches(11.7), Inches(0.5),
             f'"{d["quote"]}"', size=18, color=SECONDARY, font=FONT_B)

def t_comparison(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.0),
         d.get("title",""), size=34, bold=True, color=PRIMARY, font=FONT_H)
    # Support both flat (left_title/left_items) and nested (left: {title, bullets}) schemas
    left  = d.get("left")  or {}
    right = d.get("right") or {}
    left_title  = d.get("left_title")  or left.get("title", "")
    right_title = d.get("right_title") or right.get("title", "")
    left_items  = d.get("left_items")  or left.get("bullets")  or left.get("items")  or []
    right_items = d.get("right_items") or right.get("bullets") or right.get("items") or []
    col_w = Inches(5.6); gap = Inches(0.5)
    left_x = Inches(0.8); right_x = left_x + col_w + gap
    col_top = Inches(2.2); col_h = Inches(4.6)
    for x, title, items, accent in [
        (left_x,  left_title,  left_items,  SECONDARY),
        (right_x, right_title, right_items, PRIMARY),
    ]:
        rect(s, x, col_top, col_w, col_h, fill=PANEL, line=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.7),
             title, size=22, bold=True, color=accent, font=FONT_H)
        bullet_list(s, x + Inches(0.3), col_top + Inches(1.05),
                    col_w - Inches(0.6), col_h - Inches(1.3),
                    items, size=18, gap=Pt(8))
    # VS circle in the middle
    vs_size = Inches(0.7)
    cx = left_x + col_w + (gap - vs_size)/2
    cy = col_top + (col_h - vs_size)/2
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, vs_size, vs_size)
    circ.fill.solid(); circ.fill.fore_color.rgb = SECONDARY
    circ.line.fill.background()
    tf = circ.text_frame; tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "VS"
    r.font.name = FONT_H; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = PRIMARY

def t_diagram(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.0),
         d.get("title",""), size=34, bold=True, color=PRIMARY, font=FONT_H)
    png = d.get("mermaid_png")
    area_l, area_t = Inches(0.8), Inches(2.0)
    area_w, area_h = Inches(11.7), Inches(4.7)
    if png and os.path.exists(png):
        fit_picture(s, png, area_l, area_t, area_w, area_h)
    else:
        rect(s, area_l, area_t, area_w, area_h, fill=PANEL, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        code = d.get("mermaid_code","")
        tb = s.shapes.add_textbox(area_l+Inches(0.2), area_t+Inches(0.2),
                                  area_w-Inches(0.4), area_h-Inches(0.4))
        tf = tb.text_frame; tf.word_wrap = True
        for i, ln in enumerate(code.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = ln or " "
            r.font.name = FONT_C; r.font.size = Pt(12); r.font.color.rgb = MUTED
    if d.get("caption"):
        text(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
             d["caption"], size=14, color=SECONDARY, font=FONT_B)

def t_code(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.9),
         d.get("title",""), size=32, bold=True, color=PRIMARY, font=FONT_H)
    if d.get("language"):
        text(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.4),
             d["language"], size=14, color=SECONDARY, font=FONT_C)
    box_l, box_t = Inches(0.8), Inches(2.25)
    box_w, box_h = Inches(11.7), Inches(4.3)
    rect(s, box_l, box_t, box_w, box_h, fill=PANEL, line=BORDER,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = s.shapes.add_textbox(box_l + Inches(0.3), box_t + Inches(0.2),
                              box_w - Inches(0.6), box_h - Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate((d.get("code","") or "").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln or " "
        r.font.name = FONT_C; r.font.size = Pt(16); r.font.color.rgb = RGBColor(0xE6,0xE6,0xE6)
    if d.get("note"):
        text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
             d["note"], size=14, color=SECONDARY, font=FONT_B)

def t_takeaways(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.0),
         d.get("title","Para levar pra casa"), size=44, bold=True, color=PRIMARY, font=FONT_H)
    items = d.get("items") or []
    y = Inches(2.2); line_h = Inches(0.75)
    for i, item in enumerate(items):
        # numbered chip
        chip_w = Inches(0.55); chip_h = Inches(0.55)
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + i*line_h, chip_w, chip_h)
        chip.fill.solid(); chip.fill.fore_color.rgb = SECONDARY; chip.line.fill.background()
        tf = chip.text_frame; tf.margin_left=tf.margin_right=Emu(0)
        tf.margin_top=tf.margin_bottom=Emu(0); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i+1)
        r.font.name = FONT_H; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = PRIMARY
        text(s, Inches(1.55), y + i*line_h, Inches(11.0), line_h,
             item, size=22, color=TEXT, font=FONT_B, anchor=MSO_ANCHOR.MIDDLE)

def t_transition(s, d):
    add_bg(s)
    text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.6),
         "PRÓXIMO BLOCO", size=18, color=SECONDARY, font=FONT_B, bold=True)
    accent_bar(s, Inches(0.8), Inches(3.1), w=Inches(2.4), h=Inches(0.16))
    text(s, Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.2),
         d.get("title","Próximo passo"), size=56, bold=True, color=PRIMARY, font=FONT_H)
    if d.get("subtitle"):
        text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.8),
             d["subtitle"], size=30, color=TEXT, font=FONT_B)
    if d.get("hint"):
        text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.6),
             d["hint"], size=20, color=MUTED, font=FONT_B)

def t_story(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    if d.get("label"):
        text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.4),
             d["label"], size=14, bold=True, color=SECONDARY, font=FONT_B)
    text(s, Inches(0.8), Inches(1.35), Inches(11.7), Inches(1.0),
         d.get("title",""), size=38, bold=True, color=PRIMARY, font=FONT_H)
    timeline = d.get("timeline") or []
    y = Inches(2.55); row_h = Inches(0.5)
    for i, item in enumerate(timeline):
        ty = y + i*row_h
        # time chip
        text(s, Inches(0.8), ty, Inches(1.2), row_h,
             item.get("time",""), size=18, bold=True, color=SECONDARY,
             font=FONT_C, anchor=MSO_ANCHOR.MIDDLE)
        # dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.05),
                                 ty + Inches(0.18), Inches(0.16), Inches(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = SECONDARY
        dot.line.fill.background()
        # connector line (except last)
        if i < len(timeline) - 1:
            ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.115),
                                    ty + Inches(0.34), Inches(0.03), row_h - Inches(0.16))
            ln.fill.solid(); ln.fill.fore_color.rgb = BORDER
            ln.line.fill.background()
        text(s, Inches(2.45), ty, Inches(10.0), row_h,
             item.get("text",""), size=18, color=TEXT, font=FONT_B,
             anchor=MSO_ANCHOR.MIDDLE)
    if d.get("lesson"):
        text(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.5),
             f"→ {d['lesson']}", size=18, bold=True, color=SECONDARY, font=FONT_B)

def t_metrics(s, d):
    add_bg(s)
    accent_bar(s, Inches(0.8), Inches(0.8))
    text(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.0),
         d.get("title",""), size=38, bold=True, color=PRIMARY, font=FONT_H)
    if d.get("subtitle"):
        text(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.5),
             d["subtitle"], size=18, color=SECONDARY, font=FONT_B)
    stats = d.get("stats") or []
    # 2x2 grid
    grid_top = Inches(2.5); grid_left = Inches(0.8)
    card_w = Inches(5.85); card_h = Inches(2.0); gap = Inches(0.2)
    for i, st in enumerate(stats[:4]):
        col = i % 2; row = i // 2
        x = grid_left + col*(card_w + gap)
        y = grid_top + row*(card_h + gap)
        rect(s, x, y, card_w, card_h, fill=PANEL, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.6), Inches(0.4),
             st.get("label",""), size=14, bold=True, color=MUTED, font=FONT_B)
        # before/after
        text(s, x + Inches(0.3), y + Inches(0.65), card_w/2 - Inches(0.4), Inches(1.0),
             st.get("before",""), size=28, color=MUTED, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)
        # arrow
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + card_w/2 - Inches(0.25), y + Inches(0.95),
                                 Inches(0.5), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = SECONDARY
        arr.line.fill.background()
        text(s, x + card_w/2 + Inches(0.3), y + Inches(0.65),
             card_w/2 - Inches(0.4), Inches(1.0),
             st.get("after",""), size=36, bold=True, color=PRIMARY, font=FONT_H,
             anchor=MSO_ANCHOR.MIDDLE)
    if d.get("footnote"):
        text(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
             d["footnote"], size=11, color=MUTED, font=FONT_B)

def t_closing(s, d):
    add_bg(s)
    accent_bar(s, Inches(6.07), Inches(2.5), w=Inches(1.2), h=Inches(0.16))
    text(s, Inches(0.6), Inches(2.85), Inches(12.1), Inches(1.6),
         d.get("line",""), size=54, bold=True, color=PRIMARY, font=FONT_H,
         align=PP_ALIGN.CENTER, line_spacing=1.15)
    if d.get("line2"):
        text(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(1.6),
             d["line2"], size=38, color=SECONDARY, font=FONT_H,
             align=PP_ALIGN.CENTER, line_spacing=1.15)

def t_credits(s, d):
    add_bg(s)
    accent_bar(s, Inches(6.07), Inches(2.7), w=Inches(1.2), h=Inches(0.16))
    text(s, Inches(0.8), Inches(2.95), Inches(11.7), Inches(1.5),
         d.get("title","Obrigado!"), size=72, bold=True, color=PRIMARY, font=FONT_H,
         align=PP_ALIGN.CENTER)
    if d.get("subtitle"):
        text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.8),
             d["subtitle"], size=30, color=TEXT, font=FONT_B, align=PP_ALIGN.CENTER)
    if d.get("author"):
        text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
             d["author"], size=18, color=SECONDARY, font=FONT_B, align=PP_ALIGN.CENTER)

RENDERERS = {
    "cover": t_cover, "bio": t_bio, "agenda": t_agenda, "section": t_section,
    "question": t_question, "answer": t_answer, "content": t_content,
    "comparison": t_comparison, "diagram": t_diagram, "code": t_code,
    "story": t_story, "metrics": t_metrics, "closing": t_closing,
    "takeaways": t_takeaways, "transition": t_transition, "credits": t_credits,
}

NO_FOOTER = {"cover", "credits", "closing"}

slides = sorted(doc["slides"], key=lambda x: x.get("order", 0))
total = len(slides)
print(f"Building {total} slides…")
for i, sl in enumerate(slides, 1):
    tpl = sl.get("template","content")
    s = prs.slides.add_slide(BLANK)
    RENDERERS.get(tpl, t_content)(s, sl.get("data") or {})
    if tpl not in NO_FOOTER:
        slide_number_footer(s, i, total)
    print(f"  {i:02d}/{total} [{tpl}]")

prs.save(DST)
print(f"\nSaved: {DST}")
